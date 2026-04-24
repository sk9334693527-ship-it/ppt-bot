import os
import re
import subprocess
import pdfplumber
import pytesseract
import google.generativeai as genai
from groq import Groq
import requests

from PIL import Image, ImageEnhance, ImageFilter
from pdf2image import convert_from_path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

from telegram import Update, InputFile
from telegram.ext import Application, CommandHandler, MessageHandler, filters, ContextTypes

# 🔥 Firebase
import firebase_admin
from firebase_admin import credentials, firestore
import json

firebase_json = os.getenv("FIREBASE_CREDENTIALS")
cred = credentials.Certificate(json.loads(firebase_json))
firebase_admin.initialize_app(cred)
db = firestore.client()

# 🔥 ADMIN
ADMIN_ID = int(os.getenv("ADMIN_ID"))
ADMIN_PASSWORD = os.getenv("ADMIN_PASSWORD")

# 🔐 Admin sessions
admin_sessions = set()

# ===== CONFIG =====
BOT_TOKEN = os.getenv("BOT_TOKEN")

GEMINI_KEYS = [
    os.getenv("GEMINI_API_KEY"),
    os.getenv("GEMINI_API_KEY1"),
    os.getenv("GEMINI_API_KEY2"),
    os.getenv("GEMINI_API_KEY3"),
]

GROQ_KEYS = [
    os.getenv("GROQ_API_KEY"),
    os.getenv("GROQ_API_KEY1"),
    os.getenv("GROQ_API_KEY2"),
]

AICREDITS_KEY = os.getenv("AICREDITS_API_KEY")

GEMINI_KEYS = [k for k in GEMINI_KEYS if k]
GROQ_KEYS = [k for k in GROQ_KEYS if k]

gemini_models = []
for key in GEMINI_KEYS:
    genai.configure(api_key=key)
    gemini_models.append(genai.GenerativeModel("gemini-2.5-flash"))

groq_clients = [Groq(api_key=k) for k in GROQ_KEYS]

# ===== SAVE USER =====
def save_user(user):
    try:
        db.collection("users").document(str(user.id)).set({
            "user_id": user.id,
            "username": user.username,
            "first_name": user.first_name
        }, merge=True)
    except Exception as e:
        print("Firebase error:", e)

# ===== IMAGE ENHANCE =====
def enhance_image(img):
    # Try multiple enhancements for better OCR
    img1 = img.convert("L")
    img1 = ImageEnhance.Contrast(img1).enhance(2.5)
    img1 = img1.filter(ImageFilter.SHARPEN)

    img2 = img.convert("L")
    img2 = ImageEnhance.Brightness(img2).enhance(1.2)
    img2 = ImageEnhance.Contrast(img2).enhance(3.0)
    img2 = img2.filter(ImageFilter.SHARPEN)

    return [img1, img2, img]  # Return list: enhanced, more enhanced, and original

# ===== AICREDITS =====
def generate_aicredits(prompt):
    try:
        url = "https://api.aicredits.in/v1/chat/completions"
        headers = {
            "Authorization": f"Bearer {AICREDITS_KEY}",
            "Content-Type": "application/json"
        }
        data = {
            "model": "gpt-4o-mini",
            "messages": [{"role": "user", "content": prompt}]
        }
        res = requests.post(url, headers=headers, json=data, timeout=30)
        if res.status_code == 200:
            return res.json()["choices"][0]["message"]["content"]
    except:
        pass
    return ""

# ===== AI =====
def generate_ai(prompt):
    if AICREDITS_KEY:
        res = generate_aicredits(prompt)
        if res:
            return res

    for model in gemini_models:
        try:
            res = model.generate_content(prompt)
            if res.text:
                return res.text
        except:
            continue

    for client in groq_clients:
        try:
            chat = client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model="llama3-70b-8192"
            )
            return chat.choices[0].message.content
        except:
            continue

    return ""

# ===== PROMPT =====
FIX_PROMPT = """
तुम एक हिंदी MCQ generator हो।
काम:
1. दिए गए टेक्स्ट से केवल प्रश्न निकालो
2. मात्रा की गलती सुधारो
3. प्रश्न का अर्थ मत बदलो
4. MCQ format में बदलो
FORMAT STRICT:
प्रश्न ...
A)
B)
C)
D)
कोई extra text नहीं देना।
TEXT:
"""

# ===== PPT → PDF =====
def convert_ppt_to_pdf(ppt_path):
    subprocess.run([
        "libreoffice", "--headless",
        "--convert-to", "pdf",
        "--outdir", ".", ppt_path
    ], check=True)
    return ppt_path.replace(".pptx", ".pdf")

# ===== PPT =====
async def make_ppt(update, questions):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def set_black_background(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

    def setup_tf(tf):
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    def style_question(p):
        for run in p.runs:
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(255, 255, 0)

    def style_option(p):
        for run in p.runs:
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(255, 255, 255)


    if not questions:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_black_background(slide)
        box = slide.shapes.add_textbox(Inches(3.5), Inches(3), Inches(9), Inches(1))
        tf = box.text_frame
        setup_tf(tf)
        p = tf.paragraphs[0]
        p.text = "❌ No Data"
        style_question(p)
    else:
        for q in questions:
            lines = [l.rstrip() for l in q.split("\n") if l.strip()]
            if not lines:
                continue

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_black_background(slide)

            box = slide.shapes.add_textbox(Inches(3.5), Inches(1), Inches(9), Inches(5))
            tf = box.text_frame
            tf.clear()
            setup_tf(tf)

            # Preserve original question text, remove only leading numbering if present
            question_text = lines[0]
            question_text = re.sub(r"^\d+\.\s*", "", question_text)  # Remove leading number if present

            p = tf.paragraphs[0]
            p.text = question_text
            style_question(p)

            tf.add_paragraph().text = ""

            for opt in lines[1:]:
                p = tf.add_paragraph()
                p.text = opt
                style_option(p)

    ppt_file = "output.pptx"
    prs.save(ppt_file)

    pdf_file = convert_ppt_to_pdf(ppt_file)

    with open(ppt_file, "rb") as f:
        await update.message.reply_document(InputFile(f))

    with open(pdf_file, "rb") as f:
        await update.message.reply_document(InputFile(f))

    os.remove(ppt_file)
    os.remove(pdf_file)

# ===== HANDLERS =====
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    save_user(update.effective_user)
    await update.message.reply_text("📸 Image | ✍️ Text | 📄 PDF bhejo — PPT + PDF bana dunga")

# 🔐 ADMIN LOGIN
async def admin_login(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id

    if user_id in admin_sessions:
        await update.message.reply_text("✅ Admin Panel Open\n\n/users - Users list\n/logout - Logout")
        return

    await update.message.reply_text("🔐 Password bhejo:")
    context.user_data["awaiting_admin_password"] = True

# 🔐 LOGOUT
async def admin_logout(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id

    if user_id in admin_sessions:
        admin_sessions.remove(user_id)
        await update.message.reply_text("🚪 Logout ho gaye")
    else:
        await update.message.reply_text("❌ Pehle login karo /admin")

async def handle_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    save_user(update.effective_user)

    # 🔐 Password check
    if context.user_data.get("awaiting_admin_password"):
        if update.message.text == ADMIN_PASSWORD:
            admin_sessions.add(update.effective_user.id)
            context.user_data["awaiting_admin_password"] = False
            await update.message.reply_text("✅ Login Successful\n\n/users - Users list\n/logout - Logout")
        else:
            await update.message.reply_text("❌ Wrong Password")
        return

    fixed = generate_ai(FIX_PROMPT + update.message.text)
    if not fixed:
        await update.message.reply_text("❌ AI fail ho gaya")
        return

    questions = re.split(r"\n(?=प्रश्न)", fixed)
    await make_ppt(update, questions)

async def handle_image(update: Update, context: ContextTypes.DEFAULT_TYPE):
    save_user(update.effective_user)

    await update.message.reply_text("📸 Image process ho rahi hai...")
    photo = update.message.photo[-1]
    file = await photo.get_file()

    path = "img.jpg"
    await file.download_to_drive(path)

    img = enhance_image(Image.open(path))
    text = pytesseract.image_to_string(img, lang="hin+eng")
    os.remove(path)

    if not text or len(text.strip()) < 20:
        await update.message.reply_text("❌ Image se text nahi nikla")
        return

    fixed = generate_ai(FIX_PROMPT + text)
    if not fixed:
        await update.message.reply_text("❌ AI fail ho gaya")
        return

    questions = re.split(r"\n(?=प्रश्न)", fixed)
    await make_ppt(update, questions)

async def handle_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    save_user(update.effective_user)

    await update.message.reply_text("📄 PDF process ho raha hai...")
    doc = update.message.document
    file = await doc.get_file()

    path = "file.pdf"
    await file.download_to_drive(path)

    try:
        # Extract each page as image, OCR each, and make one slide per page
        with pdfplumber.open(path) as pdf:
            total_pages = len(pdf.pages)

        page_texts = []
        for i in range(1, total_pages + 1):
            images = convert_from_path(path, dpi=350, first_page=i, last_page=i)
            if not images:
                continue
            img_variants = enhance_image(images[0])
            page_text = ""
            for img in img_variants:
                t = pytesseract.image_to_string(img, lang="hin+eng")
                if t and len(t.strip()) > len(page_text.strip()):
                    page_text = t
            page_texts.append(page_text.strip())

        # Send preview to user
        preview = "\n---\n".join(page_texts)[:4000]
        if preview.strip():
            await update.message.reply_text("📝 Extracted text (preview):\n" + preview)
        else:
            await update.message.reply_text("❌ PDF se text nahi nikla (OCR bhi fail)")

        # Make one slide per page's text (no AI, just raw text)
        if any(page_texts):
            # Use a custom PPT function for raw text
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            def set_black_background(slide):
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 0, 0)

            for idx, txt in enumerate(page_texts, 1):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                set_black_background(slide)
                box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6.5))
                tf = box.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                p = tf.paragraphs[0]
                p.text = txt if txt else "(No text found)"
                for run in p.runs:
                    run.font.size = Pt(20)
                    run.font.color.rgb = RGBColor(255, 255, 255)

            ppt_file = "output_raw.pptx"
            prs.save(ppt_file)
            with open(ppt_file, "rb") as f:
                await update.message.reply_document(InputFile(f))
            os.remove(ppt_file)
        else:
            await update.message.reply_text("❌ Koi bhi page se text nahi nikla.")

    except Exception as e:
        await update.message.reply_text(f"❌ ERROR: {str(e)}")

    finally:
        if os.path.exists(path):
            os.remove(path)

# 🔥 ADMIN USERS
async def admin_users(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id not in admin_sessions:
        await update.message.reply_text("❌ Admin login karo /admin")
        return

    users = db.collection("users").stream()

    msg = "👥 Users List:\n\n"

    for u in users:
        data = u.to_dict()
        msg += f"ID: {data.get('user_id')}\n"
        msg += f"Username: {data.get('username')}\n\n"

    await update.message.reply_text(msg[:4000])

# ===== MAIN =====
def main():
    app = Application.builder().token(BOT_TOKEN).build()

    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("admin", admin_login))
    app.add_handler(CommandHandler("logout", admin_logout))
    app.add_handler(CommandHandler("users", admin_users))

    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))
    app.add_handler(MessageHandler(filters.PHOTO, handle_image))
    app.add_handler(MessageHandler(filters.Document.PDF, handle_pdf))

    print("🚀 Bot running...")
    app.run_polling()

if __name__ == "__main__":
    main()
